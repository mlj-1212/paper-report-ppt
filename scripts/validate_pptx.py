#!/usr/bin/env python3
"""
PPTX validation script for paper-report-ppt v4.0.

Performs three automated quality checks on a generated PPTX file:
  1. Editability  — verifies slides contain native DrawingML objects
                    (textboxes, autoshapes, pictures), not just full-slide images
  2. Image integrity — compares sha256 of media files inside the PPTX against
                    the source image_manifest.json
  3. Outline consistency — compares slide titles against outline.md page sequence

Usage:
  python validate_pptx.py <pptx_path>
  python validate_pptx.py <pptx_path> --manifest image_manifest.json --outline outline.md
  python validate_pptx.py <pptx_path> --json   # machine-readable output
"""

import argparse
import hashlib
import json
import re
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ═══════════════════════════════════════════════════════════
#  Check 1: Editability
# ═══════════════════════════════════════════════════════════

def check_editability(pptx_path):
    """Verify that slides contain native editable objects, not just full-slide images."""
    prs = Presentation(str(pptx_path))
    total_slides = len(prs.slides)
    total_shapes = 0
    total_textboxes = 0
    total_autoshapes = 0
    total_pictures = 0
    total_tables = 0
    slides_with_only_image = []

    for idx, slide in enumerate(prs.slides, 1):
        slide_shapes = list(slide.shapes)
        shape_count = len(slide_shapes)
        total_shapes += shape_count

        pic_count = 0
        non_pic_count = 0

        for shape in slide_shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                total_pictures += 1
                pic_count += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                total_autoshapes += 1
                non_pic_count += 1
            elif shape.has_text_frame:
                total_textboxes += 1
                non_pic_count += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                total_tables += 1
                non_pic_count += 1
            else:
                non_pic_count += 1

        # Flag slides that contain only a single picture (likely non-editable)
        if shape_count > 0 and pic_count > 0 and non_pic_count == 0:
            slides_with_only_image.append(idx)

    passed = total_shapes > 0 and len(slides_with_only_image) == 0

    return {
        "passed": passed,
        "total_slides": total_slides,
        "total_shapes": total_shapes,
        "textboxes": total_textboxes,
        "autoshapes": total_autoshapes,
        "pictures": total_pictures,
        "tables": total_tables,
        "slides_with_only_image": slides_with_only_image,
        "message": (
            "All slides contain native editable objects"
            if passed
            else f"{len(slides_with_only_image)} slide(s) contain only images (not editable)"
        ),
    }


# ═══════════════════════════════════════════════════════════
#  Check 2: Image integrity
# ═══════════════════════════════════════════════════════════

def _extract_media_sha256(pptx_path):
    """Extract all media files from PPTX and compute their sha256."""
    media_hashes = {}
    with zipfile.ZipFile(str(pptx_path), "r") as zf:
        for name in zf.namelist():
            if name.startswith("ppt/media/"):
                data = zf.read(name)
                sha = hashlib.sha256(data).hexdigest()
                media_hashes[name] = sha
    return media_hashes


def _load_manifest_hashes(manifest_path):
    """Load sha256 values from image_manifest.json."""
    if not manifest_path or not Path(manifest_path).exists():
        return None

    with open(manifest_path, "r", encoding="utf-8") as f:
        manifest = json.load(f)

    hashes = {}
    items = manifest if isinstance(manifest, list) else manifest.get("items", [])
    for item in items:
        sha = item.get("sha256") or item.get("source_sha256")
        filename = item.get("filename", "")
        if sha and filename:
            hashes[filename] = sha
    return hashes


def check_image_integrity(pptx_path, manifest_path=None):
    """Compare media file hashes in PPTX against source manifest."""
    media_hashes = _extract_media_sha256(pptx_path)
    manifest_hashes = _load_manifest_hashes(manifest_path)

    if not manifest_hashes:
        return {
            "passed": True,
            "media_files_found": len(media_hashes),
            "matched": 0,
            "mismatched": 0,
            "message": "No manifest provided — image integrity check skipped",
        }

    matched = 0
    mismatched = 0
    mismatches = []

    for media_name, media_sha in media_hashes.items():
        # Try to match by filename (media name may differ from manifest filename)
        media_basename = Path(media_name).name
        found = False
        for manifest_name, manifest_sha in manifest_hashes.items():
            if manifest_name in media_basename or media_basename in manifest_name:
                found = True
                if media_sha == manifest_sha:
                    matched += 1
                else:
                    mismatched += 1
                    mismatches.append({
                        "file": media_name,
                        "pptx_sha256": media_sha[:16] + "...",
                        "manifest_sha256": manifest_sha[:16] + "...",
                    })
                break

        if not found:
            # Media file not in manifest — may be a decorative shape image
            pass

    passed = mismatched == 0

    return {
        "passed": passed,
        "media_files_found": len(media_hashes),
        "matched": matched,
        "mismatched": mismatched,
        "mismatches": mismatches,
        "message": (
            f"All {matched} media files match source"
            if passed
            else f"{mismatched} media file(s) have hash mismatch"
        ),
    }


# ═══════════════════════════════════════════════════════════
#  Check 3: Outline consistency
# ═══════════════════════════════════════════════════════════

def _extract_slide_titles(pptx_path):
    """Extract title text from each slide (first substantial text frame)."""
    prs = Presentation(str(pptx_path))
    titles = []
    for slide in prs.slides:
        title_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) > 2:
                    title_text = text
                    break
        titles.append(title_text)
    return titles


def _parse_outline_titles(outline_path):
    """Extract expected page titles from outline.md."""
    if not outline_path or not Path(outline_path).exists():
        return None

    with open(outline_path, "r", encoding="utf-8") as f:
        content = f.read()

    # Look for patterns like "## P03 — Title" or "### P03: Title" or "- P03: Title"
    titles = []
    patterns = [
        r"^#+\s*(P\d+)\s*[—:\-]\s*(.+)$",
        r"^-\s*(P\d+)\s*[—:\-]\s*(.+)$",
        r"^\*\*(P\d+)\*\*\s*[—:\-]\s*(.+)$",
    ]
    for pattern in patterns:
        matches = re.findall(pattern, content, re.MULTILINE)
        if matches:
            for page_num, title in matches:
                titles.append(title.strip())
            break

    return titles if titles else None


def check_outline_consistency(pptx_path, outline_path=None):
    """Compare slide titles against outline.md page sequence."""
    slide_titles = _extract_slide_titles(pptx_path)
    outline_titles = _parse_outline_titles(outline_path)

    if not outline_titles:
        return {
            "passed": True,
            "slide_count": len(slide_titles),
            "outline_count": 0,
            "matched": 0,
            "message": "No outline provided — outline consistency check skipped",
        }

    slide_count = len(slide_titles)
    outline_count = len(outline_titles)
    matched = 0

    for i in range(min(slide_count, outline_count)):
        # Fuzzy match: check if outline title appears in slide title or vice versa
        s_title = slide_titles[i].lower().strip()
        o_title = outline_titles[i].lower().strip()
        if o_title in s_title or s_title in o_title:
            matched += 1

    passed = matched >= min(slide_count, outline_count) * 0.7  # 70% threshold

    return {
        "passed": passed,
        "slide_count": slide_count,
        "outline_count": outline_count,
        "matched": matched,
        "message": (
            f"{matched}/{min(slide_count, outline_count)} titles match outline"
            if passed
            else f"Only {matched}/{min(slide_count, outline_count)} titles match — outline inconsistency detected"
        ),
    }


# ═══════════════════════════════════════════════════════════
#  Main
# ═══════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description="Validate a generated PPTX file for editability, image integrity, and outline consistency."
    )
    parser.add_argument("pptx_path", help="Path to the PPTX file to validate")
    parser.add_argument("--manifest", "-m", help="Path to image_manifest.json")
    parser.add_argument("--outline", "-o", help="Path to outline.md")
    parser.add_argument("--json", action="store_true", help="Output JSON report")
    args = parser.parse_args()

    pptx_path = Path(args.pptx_path)
    if not pptx_path.exists():
        print(f"[ERROR] PPTX file not found: {pptx_path}", file=sys.stderr)
        return 1

    # Run all three checks
    report = {
        "pptx_path": str(pptx_path),
        "editability": check_editability(pptx_path),
        "image_integrity": check_image_integrity(pptx_path, args.manifest),
        "outline_consistency": check_outline_consistency(pptx_path, args.outline),
    }

    # Overall pass/fail
    report["overall_passed"] = all(
        report[check]["passed"] for check in ["editability", "image_integrity", "outline_consistency"]
    )

    if args.json:
        print(json.dumps(report, ensure_ascii=False, indent=2))
    else:
        print("=" * 60)
        print("PPTX 质检报告")
        print("=" * 60)
        print(f"\n文件: {pptx_path}")

        # Editability
        ed = report["editability"]
        status = "✅ 通过" if ed["passed"] else "❌ 未通过"
        print(f"\n--- 1. 可编辑性 {status} ---")
        print(f"  幻灯片数: {ed['total_slides']}")
        print(f"  总形状数: {ed['total_shapes']}")
        print(f"  文本框: {ed['textboxes']} | 形状: {ed['autoshapes']} | 图片: {ed['pictures']} | 表格: {ed['tables']}")
        if ed["slides_with_only_image"]:
            print(f"  ⚠️ 仅含图片的页: {ed['slides_with_only_image']}")

        # Image integrity
        img = report["image_integrity"]
        status = "✅ 通过" if img["passed"] else "❌ 未通过"
        print(f"\n--- 2. 图片完整性 {status} ---")
        print(f"  PPTX内媒体文件: {img['media_files_found']}")
        print(f"  匹配: {img['matched']} | 不匹配: {img['mismatched']}")
        if img.get("mismatches"):
            for m in img["mismatches"]:
                print(f"    {m['file']}: sha256 不一致")

        # Outline consistency
        ol = report["outline_consistency"]
        status = "✅ 通过" if ol["passed"] else "❌ 未通过"
        print(f"\n--- 3. 脉络一致性 {status} ---")
        print(f"  幻灯片标题数: {ol['slide_count']}")
        print(f"  大纲页数: {ol['outline_count']}")
        print(f"  匹配: {ol['matched']}")

        print(f"\n{'=' * 60}")
        if report["overall_passed"]:
            print("✅ 所有质检项通过")
        else:
            failed = [k for k in ["editability", "image_integrity", "outline_consistency"] if not report[k]["passed"]]
            print(f"❌ 质检未通过: {', '.join(failed)}")
        print("=" * 60)

    return 0 if report["overall_passed"] else 1


if __name__ == "__main__":
    sys.exit(main())
