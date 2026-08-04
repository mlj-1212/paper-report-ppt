#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
gen_pptx.py —— paper-report-ppt 技能的"直接生成"路径
==========================================
直接使用 python-pptx 生成高质量、完全可编辑的 PPTX 文件，
绕过 SVG 管线，输出质量对标 SVG 管线（每页 15-35 个形状、丰富装饰、精细布局）。

用法:
    python gen_pptx.py --input slides.json --images-dir ./images --output output.pptx --theme ref
"""

import argparse
import json
import os
import sys
from pathlib import Path

# ── 依赖检查 ──────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("错误: 缺少 python-pptx 库。请执行: pip install python-pptx")
    sys.exit(1)

# ═══════════════════════════════════════════════════════════
#  主题配置
# ═══════════════════════════════════════════════════════════

THEMES = {
    "academic": {
        # 颜色
        "background": RGBColor(0xF8, 0xF9, 0xFA),       # #F8F9FA 浅灰背景
        "header_bar": RGBColor(0x2C, 0x3E, 0x50),        # #2C3E50 深蓝灰顶部装饰条
        "accent": RGBColor(0x34, 0x98, 0xDB),             # #3498DB 蓝色强调
        "highlight": RGBColor(0xE8, 0xF4, 0xFD),         # #E8F4FD 浅蓝高亮
        "body_text": RGBColor(0x2C, 0x3E, 0x50),          # #2C3E50 正文
        "secondary_text": RGBColor(0x7F, 0x8C, 0x8D),     # #7F8C8D 次要文字
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "border_gray": RGBColor(0xBD, 0xC3, 0xC7),        # #BDC3C7 图片边框
        "decor_corner": RGBColor(0xE8, 0xF4, 0xFD),       # 角落装饰（50%透明度通过代码控制）
        "section_bg": RGBColor(0x34, 0x98, 0xDB),          # 章节分隔页背景带
        "title_band": RGBColor(0x1A, 0x52, 0x76),          # #1A5276 深蓝色标题带
        "conclusion_bg": RGBColor(0xEB, 0xF5, 0xFB),       # #EBF5FB 浅蓝色结论框背景
        "keyword_red": RGBColor(0xC0, 0x39, 0x2B),         # #C0392B 关键词红色
        "sub_title_bg": RGBColor(0xF0, 0xF7, 0xFC),        # #F0F7FC 小标题浅蓝背景
        # 字体
        "font_title": "Microsoft YaHei",
        "font_body": "Microsoft YaHei",
        "font_caption": "Microsoft YaHei",
        "title_size_pt": 28,
        "body_size_pt": 14,
        "caption_size_pt": 12,
        # 布局
        "slide_width_inch": 13.33,
        "slide_height_inch": 7.5,
        "top_bar_height_inch": 0.3,
        "header_area_height_inch": 1.1,
        "margin_inch": 0.8,
        # 装饰开关
        "decorations": False,
    },
    "minimal": {
        "background": RGBColor(0xFF, 0xFF, 0xFF),
        "header_bar": RGBColor(0x2C, 0x3E, 0x50),
        "accent": RGBColor(0x34, 0x98, 0xDB),
        "highlight": RGBColor(0xF5, 0xF5, 0xF5),
        "body_text": RGBColor(0x33, 0x33, 0x33),
        "secondary_text": RGBColor(0x99, 0x99, 0x99),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "border_gray": RGBColor(0xDD, 0xDD, 0xDD),
        "decor_corner": RGBColor(0xF0, 0xF0, 0xF0),
        "section_bg": RGBColor(0x34, 0x98, 0xDB),
        "font_title": "Microsoft YaHei",
        "font_body": "Microsoft YaHei",
        "font_caption": "Microsoft YaHei",
        "title_size_pt": 26,
        "body_size_pt": 14,
        "caption_size_pt": 12,
        "slide_width_inch": 13.33,
        "slide_height_inch": 7.5,
        "top_bar_height_inch": 0.15,
        "header_area_height_inch": 1.0,
        "margin_inch": 1.0,
        "decorations": False,
    },
    "trae": {
        # 颜色 — 自然简洁、无装饰
        "background": RGBColor(0xFF, 0xFF, 0xFF),       # #FFFFFF 纯白背景
        "header_bar": RGBColor(0x1A, 0x1A, 0x2E),        # #1A1A2E 深色顶部条
        "accent": RGBColor(0x00, 0x7A, 0xCC),             # #007ACC TRAE蓝
        "highlight": RGBColor(0xF0, 0xF6, 0xFC),         # #F0F6FC 浅蓝高亮
        "body_text": RGBColor(0x1A, 0x1A, 0x2E),          # #1A1A2E 正文
        "secondary_text": RGBColor(0x6C, 0x75, 0x7D),     # #6C757D 次要文字
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "border_gray": RGBColor(0xDE, 0xE2, 0xE6),        # #DEE2E6 图片边框
        "decor_corner": RGBColor(0xF0, 0xF6, 0xFC),       # 角落装饰（50%透明度）
        "section_bg": RGBColor(0x00, 0x7A, 0xCC),          # #007ACC 章节分隔页背景带
        "title_band": RGBColor(0x1A, 0x1A, 0x2E),          # #1A1A2E 深色标题带
        "conclusion_bg": RGBColor(0xF0, 0xF6, 0xFC),       # #F0F6FC 浅蓝结论框背景
        "keyword_red": RGBColor(0xC0, 0x39, 0x2B),         # #C0392B 关键词红色
        "sub_title_bg": RGBColor(0xF8, 0xF9, 0xFA),        # #F8F9FA 小标题浅背景
        # 字体
        "font_title": "Microsoft YaHei",
        "font_body": "Microsoft YaHei",
        "font_caption": "Microsoft YaHei",
        "title_size_pt": 28,
        "body_size_pt": 14,
        "caption_size_pt": 12,
        # 布局
        "slide_width_inch": 13.33,
        "slide_height_inch": 7.5,
        "top_bar_height_inch": 0.25,
        "header_area_height_inch": 1.1,
        "margin_inch": 0.9,
        # 装饰开关 — 无装饰
        "decorations": False,
    },
    "ref": {
        # 颜色 — 对标参考 PPT 风格（浅灰底 + 海军蓝导航 + 白色直角卡片）
        "background": RGBColor(0xEE, 0xF2, 0xF5),       # #EEF2F5 浅灰背景
        "header_bar": RGBColor(0x30, 0x43, 0x71),        # #304371 深蓝标题带
        "accent": RGBColor(0x30, 0x43, 0x71),             # #304371 深蓝强调
        "accent2": RGBColor(0x24, 0x32, 0x55),            # #243255 海军蓝（原橙色已弃用）
        "accent3": RGBColor(0x24, 0x32, 0x55),            # #243255 海军蓝
        "highlight": RGBColor(0xFF, 0xFF, 0xFF),          # #FFFFFF 白色高亮
        "body_text": RGBColor(0x33, 0x33, 0x33),          # #333333 正文
        "secondary_text": RGBColor(0x66, 0x66, 0x66),     # #666666 次要文字
        "tertiary_text": RGBColor(0x80, 0x80, 0x80),      # #808080 三级文字
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "border_gray": RGBColor(0xDD, 0xDD, 0xDD),        # #DDDDDD 边框灰
        "conclusion_bg": RGBColor(0xFF, 0xFF, 0xFF),      # #FFFFFF 结论框背景（透明感）
        "keyword_red": RGBColor(0xC0, 0x39, 0x2B),         # #C0392B 关键词红色
        "sub_title_bg": RGBColor(0xFF, 0xFF, 0xFF),        # #FFFFFF 小标题背景
        "title_band": RGBColor(0x30, 0x43, 0x71),          # #304371 标题带
        "card_bg": RGBColor(0xFF, 0xFF, 0xFF),             # #FFFFFF 卡片背景
        "nav_bg": RGBColor(0x24, 0x32, 0x55),              # #243255 导航栏背景
        "nav_active": RGBColor(0x30, 0x43, 0x71),          # #304371 激活标签
        # 字体
        "font_title": "Microsoft YaHei",
        "font_body": "Arial",
        "font_caption": "Arial",
        "title_size_pt": 21,
        "body_size_pt": 10.5,
        "caption_size_pt": 10,
        # 布局
        "slide_width_inch": 13.33,
        "slide_height_inch": 7.5,
        "top_bar_height_inch": 0.458,
        "bottom_bar_height_inch": 0.417,
        "margin_inch": 0.625,
        # 装饰开关 — 无装饰
        "decorations": False,
    },
}

# ═══════════════════════════════════════════════════════════
#  辅助函数
# ═══════════════════════════════════════════════════════════

def _set_font(run, font_name, size_pt, bold=False, color=None, italic=False):
    """设置文本 run 的字体属性，同时处理 CJK 字体兼容。"""
    run.font.name = font_name
    if bold:
        run.font.bold = True
    if italic:
        run.font.italic = True
    if size_pt:
        run.font.size = Pt(size_pt)
    if color:
        run.font.color.rgb = color
    # CJK 字体: 通过 XML 设置 eastAsia 字体名，确保中文显示正确
    rPr = run._r.get_or_add_rPr()
    nsmap = rPr.nsmap if hasattr(rPr, "nsmap") else None
    # 构造 eastAsia 属性
    ea_elem = rPr.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}eastAsia",
        {"typeface": font_name},
    )
    rPr.append(ea_elem)


def _add_textbox(slide, left, top, width, height, text, font_name, size_pt,
                color, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                word_wrap=True):
    """在幻灯片上添加一个文本框并返回 shape 对象。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, font_name, size_pt, bold=bold, color=color, italic=italic)
    return txBox


def _add_rect(slide, left, top, width, height, fill_color, line_color=None,
              line_width=None):
    """添加矩形装饰形状。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def _add_rounded_rect(slide, left, top, width, height, fill_color,
                      line_color=None, line_width=None):
    """添加圆角矩形装饰形状。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def _add_oval(slide, left, top, width, height, fill_color, line_color=None):
    """添加圆形/椭圆装饰形状。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def _set_shape_opacity(shape, opacity_pct):
    """通过 XML 设置形状填充的透明度 (0-100, 0=完全透明)。"""
    try:
        spPr = shape._element.spPr
        solidFill = spPr.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
        if solidFill is not None:
            srgbClr = solidFill.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
            if srgbClr is not None:
                ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
                alpha_elem = srgbClr.makeelement(
                    f"{{{ns}}}alpha",
                    {"val": str(int((1 - opacity_pct / 100) * 100000))},
                )
                # 移除已有的 alpha 元素
                for old in srgbClr.findall(f"{{{ns}}}alpha"):
                    srgbClr.remove(old)
                srgbClr.append(alpha_elem)
    except Exception:
        pass  # 透明度设置失败不影响主流程


def _add_notes(slide, notes_text):
    """为幻灯片添加演讲者备注。"""
    if notes_text:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes_text


# ═══════════════════════════════════════════════════════════
#  SlideBuilder 主类
# ═══════════════════════════════════════════════════════════

class SlideBuilder:
    """PPTX 幻灯片构建器，支持多种页面类型和装饰。"""

    def __init__(self, theme_name="academic"):
        self.theme_name = theme_name
        self.theme = THEMES.get(theme_name, THEMES["academic"])
        self.sw = Inches(self.theme["slide_width_inch"])
        self.sh = Inches(self.theme["slide_height_inch"])

    def create_presentation(self):
        """创建空白演示文稿，设置 16:9 尺寸。"""
        prs = Presentation()
        prs.slide_width = self.sw
        prs.slide_height = self.sh
        return prs

    # ── 公共装饰方法 ──────────────────────────────────────

    def _add_top_bar(self, slide):
        """每页顶部的装饰横条 (深蓝灰)。"""
        _add_rect(
            slide,
            left=Emu(0), top=Emu(0),
            width=self.sw, height=Inches(self.theme["top_bar_height_inch"]),
            fill_color=self.theme["header_bar"],
        )

    def _add_bottom_bar(self, slide, page_num):
        """底部页码条。"""
        # 底部细线
        _add_rect(
            slide,
            left=Emu(0),
            top=self.sh - Inches(0.05),
            width=self.sw,
            height=Inches(0.05),
            fill_color=self.theme["accent"],
        )
        # 页码
        _add_textbox(
            slide,
            left=self.sw - Inches(1.0),
            top=self.sh - Inches(0.45),
            width=Inches(0.8),
            height=Inches(0.35),
            text=str(page_num),
            font_name=self.theme["font_caption"],
            size_pt=10,
            color=self.theme["secondary_text"],
            alignment=PP_ALIGN.RIGHT,
        )

    def _add_corner_decorations(self, slide):
        """在角落添加半透明装饰形状（仅 academic 主题）。"""
        if not self.theme["decorations"]:
            return
        t = self.theme
        # 右上角小圆形
        c1 = _add_oval(
            slide,
            left=self.sw - Inches(0.8),
            top=Inches(0.5),
            width=Inches(0.5),
            height=Inches(0.5),
            fill_color=t["decor_corner"],
        )
        _set_shape_opacity(c1, 50)
        # 左下角小方形
        s1 = _add_rect(
            slide,
            left=Inches(0.3),
            top=self.sh - Inches(1.0),
            width=Inches(0.35),
            height=Inches(0.35),
            fill_color=t["decor_corner"],
        )
        _set_shape_opacity(s1, 50)
        # 右下角圆形
        c2 = _add_oval(
            slide,
            left=self.sw - Inches(1.5),
            top=self.sh - Inches(1.3),
            width=Inches(0.3),
            height=Inches(0.3),
            fill_color=t["accent"],
        )
        _set_shape_opacity(c2, 70)

    # ── TRAE 风格辅助方法 ────────────────────────────────

    # 学术关键词列表（用于 _add_rich_textbox 智能高亮）
    ACADEMIC_KEYWORDS = [
        "增加", "减少", "升高", "降低", "促进", "抑制", "上调", "下调",
        "增强", "减弱", "正向调控", "负向调控", "加重", "减轻",
        "稳定", "降解", "激活", "破坏",
    ]

    def _add_title_header(self, slide, title_text, theme, kicker=None):
        """在页面顶部画一个全宽深蓝色矩形标题带（替代原 _add_top_bar + 单独标题）。

        - 高 0.8 英寸，颜色用 title_band
        - 标题文字白色、加粗、左对齐（左边距 0.5 英寸，垂直居中）
        - 若提供 kicker，在右上角画一个浅色小标签（圆角矩形，白字）
        """
        band_height = Inches(0.8)
        # 全宽深蓝色矩形条
        _add_rect(
            slide,
            left=Emu(0), top=Emu(0),
            width=self.sw, height=band_height,
            fill_color=theme["title_band"],
        )
        # 标题文字（白色、加粗、左对齐、垂直居中）
        title_box = _add_textbox(
            slide,
            left=Inches(0.5),
            top=Emu(0),
            width=self.sw - Inches(1.0),
            height=band_height,
            text=title_text,
            font_name=theme["font_title"],
            size_pt=theme["title_size_pt"],
            color=theme["white"],
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )
        # 垂直居中
        try:
            title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass

        # 右上角 kicker 小标签
        if kicker:
            kicker_w = Inches(1.8)
            kicker_h = Inches(0.35)
            kicker_left = self.sw - kicker_w - Inches(0.4)
            kicker_top = (band_height - kicker_h) // 2
            k_shape = _add_rounded_rect(
                slide,
                left=kicker_left,
                top=kicker_top,
                width=kicker_w,
                height=kicker_h,
                fill_color=theme["accent"],
            )
            k_tf = k_shape.text_frame
            k_tf.word_wrap = False
            try:
                k_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
            k_p = k_tf.paragraphs[0]
            k_p.alignment = PP_ALIGN.CENTER
            k_run = k_p.add_run()
            k_run.text = kicker
            _set_font(k_run, theme["font_body"], 11, bold=True, color=theme["white"])

    def _add_conclusion_box(self, slide, left, top, width, text, theme):
        """画一个浅蓝色圆角矩形结论框，左侧带蓝色竖线装饰。

        - conclusion_bg 填充，无边框
        - 左侧 accent 色竖线（宽 0.06 英寸）
        - 文字为 header_bar 色（深蓝），加粗，size 13pt，前加"结论："前缀（加粗）
        """
        box_height = Inches(0.7)
        # 浅蓝色圆角矩形背景
        _add_rounded_rect(
            slide,
            left=left,
            top=top,
            width=width,
            height=box_height,
            fill_color=theme["conclusion_bg"],
            line_color=None,
        )
        # 左侧蓝色竖线装饰
        _add_rect(
            slide,
            left=left,
            top=top,
            width=Inches(0.06),
            height=box_height,
            fill_color=theme["accent"],
        )
        # 文字框（"结论："前缀加粗 + 正文）
        txBox = slide.shapes.add_textbox(
            left + Inches(0.25),
            top,
            width - Inches(0.4),
            box_height,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        # "结论：" 前缀
        prefix_run = p.add_run()
        prefix_run.text = "结论："
        _set_font(prefix_run, theme["font_body"], 13, bold=True,
                  color=theme["header_bar"])
        # 正文（支持关键词高亮）
        self._add_keyword_runs(p, text, theme["font_body"], 13,
                               theme["header_bar"], theme)
        return box_height

    def _add_keyword_runs(self, paragraph, text, font_name, size_pt,
                          default_color, theme, bold=False):
        """将文本按学术关键词切分，添加多个 run（关键词标红加粗）。

        关键词用 keyword_red + bold，其余用 default_color。
        """
        keywords = self.ACADEMIC_KEYWORDS
        # 按出现位置切分文本
        remaining = text
        while remaining:
            # 找最早出现的关键词
            earliest_pos = -1
            earliest_kw = None
            for kw in keywords:
                pos = remaining.find(kw)
                if pos != -1 and (earliest_pos == -1 or pos < earliest_pos):
                    earliest_pos = pos
                    earliest_kw = kw
            if earliest_pos == -1:
                # 没有关键词了，输出剩余文本
                run = paragraph.add_run()
                run.text = remaining
                _set_font(run, font_name, size_pt, bold=bold, color=default_color)
                break
            # 输出关键词前的普通文本
            if earliest_pos > 0:
                normal_text = remaining[:earliest_pos]
                run = paragraph.add_run()
                run.text = normal_text
                _set_font(run, font_name, size_pt, bold=bold, color=default_color)
            # 输出关键词（标红加粗）
            kw_run = paragraph.add_run()
            kw_run.text = earliest_kw
            _set_font(kw_run, font_name, size_pt, bold=True,
                      color=theme["keyword_red"])
            # 继续处理剩余文本
            remaining = remaining[earliest_pos + len(earliest_kw):]

    def _add_rich_textbox(self, slide, left, top, width, height, text,
                          font_name, size_pt, color, theme, bold=False):
        """智能关键词高亮文本框：检测学术关键词并标红加粗。

        关键词部分用 keyword_red + bold，其余用正常 color。
        用 add_run 机制，处理一段文本中多个关键词的情况。
        """
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        self._add_keyword_runs(p, text, font_name, size_pt, color, theme, bold=bold)
        return txBox

    # ── ref 主题专用辅助方法 ────────────────────────────

    def _ref_add_nav_bar(self, slide, page_num):
        """在幻灯片顶部添加 4 个等宽导航标签。

        标签宽度 = 13.333 / 4 = 3.333 英寸，高度 0.458 英寸。
        激活标签用 nav_active 色（#304371），其余用 nav_bg 色（#243255）。
        根据 page_num 自动判断激活哪个标签。
        """
        t = self.theme
        nav_labels = ["背景", "结果", "机制", "结论"]
        # 根据 page_num 决定激活标签索引
        if page_num <= 5:
            active_idx = 0
        elif page_num <= 9:
            active_idx = 1
        elif page_num <= 13:
            active_idx = 2
        else:
            active_idx = 3

        tab_width = Inches(13.333 / 4)  # 3.333"
        tab_height = Inches(0.458)
        for i, label in enumerate(nav_labels):
            left = i * tab_width
            color = t["nav_active"] if i == active_idx else t["nav_bg"]
            _add_rect(slide, left, Emu(0), tab_width, tab_height, color)
            # 标签文字（白色加粗居中）
            label_box = _add_textbox(
                slide, left, Emu(0), tab_width, tab_height, label,
                t["font_title"], 12, t["white"], bold=True,
                alignment=PP_ALIGN.CENTER,
            )
            try:
                label_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass

    def _ref_add_bottom_bar(self, slide, page_num):
        """底部全宽条（#304371），右侧白色页码。"""
        t = self.theme
        bar_h = Inches(t["bottom_bar_height_inch"])
        _add_rect(
            slide,
            left=Emu(0),
            top=self.sh - bar_h,
            width=self.sw,
            height=bar_h,
            fill_color=t["header_bar"],
        )
        # 右下角页码（白色 9pt）
        _add_textbox(
            slide,
            left=self.sw - Inches(1.2),
            top=self.sh - bar_h,
            width=Inches(1.0),
            height=bar_h,
            text=str(page_num),
            font_name=t["font_caption"],
            size_pt=9,
            color=t["white"],
            alignment=PP_ALIGN.RIGHT,
        )

    def _ref_add_title(self, slide, title_text):
        """导航栏下方的页面标题，左对齐，下方带短下划线。

        标题位于 margin_inch（0.625"）左侧位置，下划线宽 1.25" 高 0.031"。
        返回标题区域底部 y 坐标（供后续内容定位）。
        """
        t = self.theme
        margin = Inches(t["margin_inch"])
        title_top = Inches(t["top_bar_height_inch"]) + Inches(0.12)
        # 标题文字
        _add_textbox(
            slide,
            left=margin,
            top=title_top,
            width=self.sw - margin * 2,
            height=Inches(0.5),
            text=title_text,
            font_name=t["font_title"],
            size_pt=t["title_size_pt"],
            color=t["accent3"],
            bold=True,
            alignment=PP_ALIGN.LEFT,
        )
        # 短下划线
        underline_top = title_top + Inches(0.48)
        _add_rect(
            slide,
            left=margin,
            top=underline_top,
            width=Inches(1.25),
            height=Inches(0.031),
            fill_color=t["accent"],
        )
        # 返回标题区域底部 y（下划线之下 0.1" 留白）
        return underline_top + Inches(0.13)

    def _ref_add_card(self, slide, left, top, width, height, accent_color=None):
        """添加白色直角矩形卡片（无圆角、无左侧 accent 竖条）。

        - 卡片背景：card_bg（白色），直角矩形
        返回卡片背景形状对象（便于后续在其上叠加文字时调整层级）。
        """
        t = self.theme
        # 白色直角矩形背景
        card = _add_rect(
            slide, left, top, width, height, t["card_bg"], line_color=None
        )
        return card

    def _ref_add_text_lines(self, slide, left, top, width, lines,
                            font_size, color, line_height=0.48,
                            bold=False, theme=None):
        """添加多行文本（无项目符号），每行一个独立文本框。

        支持关键词高亮（通过 _add_keyword_runs 实现）。
        返回最后一行底部的 y 坐标。
        """
        t = theme if theme is not None else self.theme
        y = top
        for line in lines:
            # 使用支持关键词高亮的文本框
            txBox = slide.shapes.add_textbox(left, y, width, Inches(line_height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            self._add_keyword_runs(
                p, line, t["font_body"], font_size, color, t, bold=bold
            )
            y += Inches(line_height)
        return y

    def _ref_add_conclusion_box(self, slide, left, top, width, text):
        """ref 主题专用结论框：白色直角矩形背景 + 加粗文字（无 accent 竖条）。

        高度根据文本自适应（简单按 1 行处理，高度 0.6"）。
        文字使用 "结论：" 前缀加粗 + 正文（支持关键词高亮）。
        返回框高度。
        """
        t = self.theme
        box_height = Inches(0.75)
        # 白色直角矩形背景
        _add_rect(
            slide,
            left=left,
            top=top,
            width=width,
            height=box_height,
            fill_color=t["conclusion_bg"],
            line_color=None,
        )
        # 文字框（"结论：" 前缀加粗 + 正文）
        txBox = slide.shapes.add_textbox(
            left + Inches(0.15),
            top,
            width - Inches(0.3),
            box_height,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        prefix_run = p.add_run()
        prefix_run.text = "结论："
        _set_font(prefix_run, t["font_body"], 10.5, bold=True,
                  color=t["accent3"])
        # 正文（支持关键词高亮）
        self._add_keyword_runs(p, text, t["font_body"], 10.5,
                               t["body_text"], t)
        return box_height

    def _ref_set_background(self, slide):
        """设置 ref 主题背景色（浅灰 #EEF2F5）。"""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self.theme["background"]

    # ── 各页面类型构建方法 ────────────────────────────────

    def build_cover(self, prs, slide_data, page_num):
        """构建封面页 —— 简洁学术风格，细线装饰，居中对齐。"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        t = self.theme

        # 背景色
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = t["background"]

        # 顶部细条装饰（全宽，高度 0.06 英寸）
        _add_rect(
            slide,
            left=Emu(0), top=Emu(0),
            width=self.sw, height=Inches(0.06),
            fill_color=t["title_band"],
        )

        # 底部细条装饰
        _add_rect(
            slide,
            left=Emu(0), top=self.sh - Inches(0.06),
            width=self.sw, height=Inches(0.06),
            fill_color=t["title_band"],
        )

        # 左侧竖线装饰（细线）
        _add_rect(
            slide,
            left=Inches(0.62), top=Inches(1.35),
            width=Inches(0.04), height=Inches(4.79),
            fill_color=t["accent"],
        )

        # 右上角 L 形小方块装饰（两段细线组成）
        _add_rect(
            slide,
            left=self.sw - Inches(1.45), top=Inches(1.35),
            width=Inches(0.83), height=Inches(0.04),
            fill_color=t["accent"],
        )
        _add_rect(
            slide,
            left=self.sw - Inches(0.66), top=Inches(1.35),
            width=Inches(0.04), height=Inches(0.83),
            fill_color=t["accent"],
        )

        # 左下角 L 形小方块装饰
        _add_rect(
            slide,
            left=Inches(0.62), top=Inches(6.10),
            width=Inches(0.04), height=Inches(0.83),
            fill_color=t["accent"],
        )
        _add_rect(
            slide,
            left=Inches(0.62), top=Inches(6.90),
            width=Inches(0.83), height=Inches(0.04),
            fill_color=t["accent"],
        )

        # 标签 "文献精读汇报"
        label = slide_data.get("label", "文献精读汇报")
        if label:
            _add_textbox(
                slide,
                left=Inches(5.0), top=Inches(1.15),
                width=Inches(3.33), height=Inches(0.38),
                text=label,
                font_name=t["font_body"],
                size_pt=14,
                color=t["secondary_text"],
                alignment=PP_ALIGN.CENTER,
            )

        # 主标题 — 居中大号加粗
        title_text = slide_data.get("title", "论文标题")
        _add_textbox(
            slide,
            left=Inches(0.5),
            top=Inches(1.97),
            width=self.sw - Inches(1.0),
            height=Inches(0.58),
            text=title_text,
            font_name=t["font_title"],
            size_pt=36,
            color=t["header_bar"],
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # 分隔线
        _add_rect(
            slide,
            left=Inches(5.10), top=Inches(2.60),
            width=Inches(3.12), height=Inches(0.01),
            fill_color=t["accent"],
        )

        # 英文副标题（如果有）
        en_subtitle = slide_data.get("en_subtitle", "")
        if en_subtitle:
            _add_textbox(
                slide,
                left=Inches(1.0),
                top=Inches(2.86),
                width=self.sw - Inches(2.0),
                height=Inches(0.58),
                text=en_subtitle,
                font_name=t["font_body"],
                size_pt=16,
                color=t["secondary_text"],
                alignment=PP_ALIGN.CENTER,
            )
        else:
            # 回退：用 subtitle 字段
            subtitle = slide_data.get("subtitle", "")
            if subtitle:
                _add_textbox(
                    slide,
                    left=Inches(2.0), top=Inches(3.8),
                    width=self.sw - Inches(4.0), height=Inches(0.8),
                    text=subtitle,
                    font_name=t["font_body"],
                    size_pt=20,
                    color=t["secondary_text"],
                    alignment=PP_ALIGN.CENTER,
                )

        # 作者 / 期刊信息框（浅蓝背景圆角矩形）
        info_lines = slide_data.get("bullets", [])
        if info_lines:
            # 背景框
            info_text = "\n".join(info_lines)
            _add_rounded_rect(
                slide,
                left=Inches(4.0), top=Inches(3.85),
                width=Inches(5.33), height=Inches(0.73),
                fill_color=t["highlight"],
                line_color=None,
            )
            _add_textbox(
                slide,
                left=Inches(4.0), top=Inches(4.01),
                width=Inches(5.33), height=Inches(0.55),
                text=info_text,
                font_name=t["font_body"],
                size_pt=14,
                color=t["body_text"],
                alignment=PP_ALIGN.CENTER,
            )

        # 底部信息（汇报人 + 日期）
        _add_textbox(
            slide,
            left=Inches(2.0), top=Inches(5.41),
            width=self.sw - Inches(4.0), height=Inches(0.39),
            text=slide_data.get("footer_text", "研究生组会汇报"),
            font_name=t["font_body"],
            size_pt=16,
            color=t["body_text"],
            alignment=PP_ALIGN.CENTER,
        )
        _add_textbox(
            slide,
            left=Inches(2.0), top=Inches(5.92),
            width=self.sw - Inches(4.0), height=Inches(0.22),
            text=slide_data.get("presenter", "汇报人：XXX"),
            font_name=t["font_body"],
            size_pt=12,
            color=t["secondary_text"],
            alignment=PP_ALIGN.CENTER,
        )
        _add_textbox(
            slide,
            left=Inches(2.0), top=Inches(6.18),
            width=self.sw - Inches(4.0), height=Inches(0.22),
            text=slide_data.get("date", "日期：2026 年 7 月"),
            font_name=t["font_body"],
            size_pt=12,
            color=t["secondary_text"],
            alignment=PP_ALIGN.CENTER,
        )

        _add_notes(slide, slide_data.get("notes", ""))

    def build_toc(self, prs, slide_data, page_num):
        """构建目录页。"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = self.theme

        # 背景色
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = t["background"]

        self._add_top_bar(slide)

        # 标题
        _add_textbox(
            slide,
            left=t.get("margin_inch", 0.8) and Inches(0.8),
            top=Inches(0.55),
            width=Inches(6.0),
            height=Inches(0.7),
            text=slide_data.get("title", "目录"),
            font_name=t["font_title"],
            size_pt=t["title_size_pt"],
            color=t["header_bar"],
            bold=True,
        )

        # 装饰左侧竖条
        _add_rect(
            slide,
            left=Inches(0.5),
            top=Inches(1.5),
            width=Inches(0.06),
            height=Inches(5.2),
            fill_color=t["accent"],
        )

        sections = slide_data.get("sections", [])
        left_col_x = Inches(1.2)
        right_col_x = Inches(7.0)
        y_start = Inches(1.8)
        line_height = Inches(0.7)

        for i, section in enumerate(sections):
            y = y_start + i * line_height
            # 编号圆圈
            circle = _add_oval(
                slide,
                left=left_col_x,
                top=y + Inches(0.05),
                width=Inches(0.45),
                height=Inches(0.45),
                fill_color=t["accent"],
            )
            # 编号文字
            circle_tf = circle.text_frame
            circle_tf.word_wrap = False
            p = circle_tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(i + 1)
            _set_font(run, t["font_body"], 14, bold=True, color=t["white"])

            # 章节名 —— 兼容字符串和字典两种 sections 格式
            section_title = _extract_section_title(section)

            _add_textbox(
                slide,
                left=left_col_x + Inches(0.65),
                top=y,
                width=Inches(4.5),
                height=Inches(0.45),
                text=section_title,
                font_name=t["font_body"],
                size_pt=16,
                color=t["body_text"],
                bold=True,
            )

            # 右侧简要描述
            highlights = slide_data.get("highlights", [])
            if i < len(highlights):
                desc = highlights[i].get("content", "")
                _add_textbox(
                    slide,
                    left=right_col_x,
                    top=y + Inches(0.05),
                    width=Inches(5.0),
                    height=Inches(0.45),
                    text=desc,
                    font_name=t["font_body"],
                    size_pt=13,
                    color=t["secondary_text"],
                )

            # 分隔线（非最后一项）
            if i < len(sections) - 1:
                _add_rect(
                    slide,
                    left=left_col_x,
                    top=y + line_height - Inches(0.1),
                    width=self.sw - Inches(2.5),
                    height=Inches(0.01),
                    fill_color=t["border_gray"],
                )

        self._add_corner_decorations(slide)
        self._add_bottom_bar(slide, page_num)
        _add_notes(slide, slide_data.get("notes", ""))

    def build_section(self, prs, slide_data, page_num):
        """构建章节分隔页。"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = self.theme

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = t["background"]

        # 顶部 title_band 色带（替代原 header_bar 细条）
        _add_rect(
            slide,
            left=Emu(0), top=Emu(0),
            width=self.sw, height=Inches(0.4),
            fill_color=t["title_band"],
        )

        # 上方装饰线
        _add_rect(
            slide,
            left=Inches(3.0),
            top=Inches(2.2),
            width=Inches(7.33),
            height=Inches(0.04),
            fill_color=t["accent"],
        )

        # 彩色强调带
        _add_rect(
            slide,
            left=Inches(3.0),
            top=Inches(2.5),
            width=Inches(7.33),
            height=Inches(0.12),
            fill_color=t["title_band"],
        )

        # 章节编号（优先用 section_no，回退到 page_num）
        section_no = slide_data.get("section_no", page_num)
        _add_textbox(
            slide,
            left=Inches(3.0),
            top=Inches(2.9),
            width=Inches(7.33),
            height=Inches(0.7),
            text=f"Section {section_no}",
            font_name=t["font_body"],
            size_pt=18,
            color=t["accent"],
            bold=False,
            alignment=PP_ALIGN.CENTER,
        )

        # 章节名 — 居中大号
        section_title = slide_data.get("title", "")
        _add_textbox(
            slide,
            left=Inches(2.0),
            top=Inches(3.5),
            width=self.sw - Inches(4.0),
            height=Inches(1.2),
            text=section_title,
            font_name=t["font_title"],
            size_pt=36,
            color=t["header_bar"],
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # 下方装饰线
        _add_rect(
            slide,
            left=Inches(3.0),
            top=Inches(5.2),
            width=Inches(7.33),
            height=Inches(0.04),
            fill_color=t["accent"],
        )

        # 装饰: 两侧小圆
        if t["decorations"]:
            c1 = _add_oval(
                slide,
                left=Inches(1.5),
                top=Inches(3.5),
                width=Inches(0.6),
                height=Inches(0.6),
                fill_color=t["decor_corner"],
            )
            _set_shape_opacity(c1, 40)
            c2 = _add_oval(
                slide,
                left=self.sw - Inches(2.3),
                top=Inches(3.5),
                width=Inches(0.6),
                height=Inches(0.6),
                fill_color=t["decor_corner"],
            )
            _set_shape_opacity(c2, 40)

        self._add_bottom_bar(slide, page_num)
        _add_notes(slide, slide_data.get("notes", ""))

    def build_content(self, prs, slide_data, page_num):
        """构建内容页（带要点列表）。"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = self.theme
        margin = Inches(t["margin_inch"])

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = t["background"]

        # 标题带（替代原 _add_top_bar + 单独标题）
        kicker = slide_data.get("kicker")
        self._add_title_header(slide, slide_data.get("title", ""), t, kicker=kicker)

        # 小标题（如果有 sub_title）
        sub_y = Inches(0.92)
        sub_title = slide_data.get("sub_title")
        if sub_title:
            # 蓝色竖线装饰
            _add_rect(
                slide,
                left=margin,
                top=sub_y + Inches(0.02),
                width=Inches(0.06),
                height=Inches(0.32),
                fill_color=t["accent"],
            )
            # 小标题浅蓝背景
            _add_rounded_rect(
                slide,
                left=margin + Inches(0.15),
                top=sub_y,
                width=self.sw - margin * 2 - Inches(0.15),
                height=Inches(0.36),
                fill_color=t["sub_title_bg"],
                line_color=None,
            )
            _add_textbox(
                slide,
                left=margin + Inches(0.3),
                top=sub_y + Inches(0.02),
                width=self.sw - margin * 2 - Inches(0.45),
                height=Inches(0.32),
                text=sub_title,
                font_name=t["font_body"],
                size_pt=14,
                color=t["title_band"],
                bold=True,
            )
            body_top = sub_y + Inches(0.52)
        else:
            body_top = Inches(1.15)

        # 要点列表
        bullets = slide_data.get("bullets", [])
        bullet_left = margin + Inches(0.3)
        bullet_width = self.sw - margin * 2 - Inches(0.6)
        bullet_height = Inches(0.55)

        for i, bullet in enumerate(bullets):
            y = body_top + i * Inches(0.7)
            num = i + 1

            # 编号圆形 — 白色底 + 蓝色边框 + 居中数字
            circle_size = Inches(0.38)
            circle_left = margin + Inches(0.05)
            circle_top = y + Inches(0.05)
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, circle_left, circle_top,
                circle_size, circle_size
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = t["white"]
            circle.line.color.rgb = t["accent"]
            circle.line.width = Pt(1.5)
            # 数字文字
            ctf = circle.text_frame
            ctf.word_wrap = False
            try:
                ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
            cp = ctf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            crun = cp.add_run()
            crun.text = str(num)
            _set_font(crun, t["font_body"], 12, bold=True, color=t["accent"])

            # 要点文本框（支持关键词高亮），左侧留出编号空间
            self._add_rich_textbox(
                slide,
                left=bullet_left + Inches(0.25),
                top=y,
                width=bullet_width - Inches(0.25),
                height=bullet_height,
                text=bullet,
                font_name=t["font_body"],
                size_pt=t["body_size_pt"],
                color=t["body_text"],
                theme=t,
            )

        # 高亮框
        highlights = slide_data.get("highlights", [])
        hl_y = body_top + len(bullets) * Inches(0.7) + Inches(0.2)
        if highlights:
            for j, hl in enumerate(highlights):
                hl_box = _add_rounded_rect(
                    slide,
                    left=margin + Inches(0.3),
                    top=hl_y + j * Inches(1.0),
                    width=Inches(5.0),
                    height=Inches(0.8),
                    fill_color=t["highlight"],
                    line_color=t["accent"],
                    line_width=1,
                )
                # 高亮标题
                _add_textbox(
                    slide,
                    left=margin + Inches(0.6),
                    top=hl_y + j * Inches(1.0) + Inches(0.05),
                    width=Inches(4.5),
                    height=Inches(0.3),
                    text=hl.get("title", ""),
                    font_name=t["font_body"],
                    size_pt=13,
                    color=t["accent"],
                    bold=True,
                )
                # 高亮内容
                _add_textbox(
                    slide,
                    left=margin + Inches(0.6),
                    top=hl_y + j * Inches(1.0) + Inches(0.35),
                    width=Inches(4.5),
                    height=Inches(0.4),
                    text=hl.get("content", ""),
                    font_name=t["font_body"],
                    size_pt=12,
                    color=t["body_text"],
                )

        # 结论框（如果有 conclusion 字段）
        conclusion = slide_data.get("conclusion")
        if conclusion:
            concl_top = self.sh - Inches(1.5)
            self._add_conclusion_box(
                slide,
                left=margin + Inches(0.3),
                top=concl_top,
                width=self.sw - margin * 2 - Inches(0.6),
                text=conclusion,
                theme=t,
            )

        self._add_corner_decorations(slide)
        self._add_bottom_bar(slide, page_num)
        _add_notes(slide, slide_data.get("notes", ""))

    def build_figure(self, prs, slide_data, page_num, images_dir=None):
        """构建图表页，支持左右分栏（文字+图片）或居中图片两种布局。"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = self.theme

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = t["background"]

        self._add_title_header(slide, slide_data.get("title", "图表"), t,
                               kicker=slide_data.get("kicker"))

        bullets = slide_data.get("bullets", [])
        has_bullets = bool(bullets)

        if has_bullets:
            # ── 左右分栏布局：左侧文字分析，右侧图片 ──
            text_col_width = Inches(5.0)
            img_col_width = self.sw - text_col_width - Inches(1.6)
            col_gap = Inches(0.4)
            text_left = Inches(0.8)
            img_left = text_left + text_col_width + col_gap
            content_top = Inches(1.15)
            content_height = self.sh - Inches(2.5)

            # 左侧文字区：极淡背景 + 无边框，靠留白自然区分
            _add_rounded_rect(
                slide,
                left=text_left - Inches(0.05),
                top=content_top - Inches(0.05),
                width=text_col_width + Inches(0.1),
                height=content_height + Inches(0.1),
                fill_color=RGBColor(0xFC, 0xFC, 0xFC),  # 极浅灰，几乎看不见
                line_color=None,
            )

            # 小标题（如果有 sub_title 或 analysis_title）
            sub_title = slide_data.get("sub_title") or slide_data.get("analysis_title")
            if sub_title:
                # 蓝色竖线装饰
                _add_rect(
                    slide,
                    left=text_left + Inches(0.05),
                    top=content_top + Inches(0.1),
                    width=Inches(0.06),
                    height=Inches(0.32),
                    fill_color=t["accent"],
                )
                # 小标题浅蓝背景
                _add_rounded_rect(
                    slide,
                    left=text_left + Inches(0.18),
                    top=content_top + Inches(0.08),
                    width=text_col_width - Inches(0.25),
                    height=Inches(0.36),
                    fill_color=t["sub_title_bg"],
                    line_color=None,
                )
                _add_textbox(
                    slide,
                    left=text_left + Inches(0.32),
                    top=content_top + Inches(0.1),
                    width=text_col_width - Inches(0.4),
                    height=Inches(0.32),
                    text=sub_title,
                    font_name=t["font_body"],
                    size_pt=14,
                    color=t["title_band"],
                    bold=True,
                )
                body_top = content_top + Inches(0.55)
            else:
                # 左侧蓝色竖线装饰（仿TRAE风格，强调文字区）
                _add_rect(
                    slide,
                    left=text_left + Inches(0.05),
                    top=content_top + Inches(0.2),
                    width=Inches(0.04),
                    height=Inches(0.6),
                    fill_color=t["accent"],
                )
                body_top = content_top + Inches(0.2)

            # 要点列表
            bullet_left = text_left + Inches(0.3)
            bullet_width = text_col_width - Inches(0.6)
            bullet_height = Inches(0.55)
            max_bullets = min(len(bullets), 6)  # 防止溢出

            for i, bullet in enumerate(bullets[:max_bullets]):
                y = body_top + i * Inches(0.7)
                # 子弹符号
                _add_oval(
                    slide,
                    left=text_left + Inches(0.08),
                    top=y + Inches(0.12),
                    width=Inches(0.15),
                    height=Inches(0.15),
                    fill_color=t["accent"],
                )
                # 要点文本（支持关键词高亮）
                self._add_rich_textbox(
                    slide,
                    left=bullet_left,
                    top=y,
                    width=bullet_width,
                    height=bullet_height,
                    text=bullet,
                    font_name=t["font_body"],
                    size_pt=t["body_size_pt"],
                    color=t["body_text"],
                    theme=t,
                )

            # 结论框（如果有 conclusion 或 key_message）
            conclusion = slide_data.get("conclusion")
            if not conclusion:
                conclusion = slide_data.get("key_message")
            if conclusion:
                concl_top = content_top + content_height - Inches(0.8)
                self._add_conclusion_box(
                    slide,
                    left=text_left + Inches(0.1),
                    top=concl_top,
                    width=text_col_width - Inches(0.2),
                    text=conclusion,
                    theme=t,
                )

            # 右侧图片区：极淡背景 + 无边框
            _add_rounded_rect(
                slide,
                left=img_left - Inches(0.05),
                top=content_top - Inches(0.05),
                width=img_col_width + Inches(0.1),
                height=content_height + Inches(0.1),
                fill_color=RGBColor(0xFA, 0xFA, 0xFA),  # 比左侧再浅一点
                line_color=None,
            )

            # 在右侧区域加载图片（为图注预留底部空间）
            img_area_height_caption = content_height - Inches(0.5)
            self._place_image_in_area(
                slide, slide_data, images_dir,
                img_left, content_top, img_col_width, img_area_height_caption, t
            )

            # 图注放在右侧图片正下方
            caption = slide_data.get("image_caption", "")
            if caption:
                _add_textbox(
                    slide,
                    left=img_left,
                    top=content_top + img_area_height_caption + Inches(0.05),
                    width=img_col_width,
                    height=Inches(0.4),
                    text=caption,
                    font_name=t["font_caption"],
                    size_pt=t["caption_size_pt"],
                    color=t["secondary_text"],
                    italic=True,
                    alignment=PP_ALIGN.CENTER,
                )

        else:
            # ── 无文字时：居中图片布局（原版） ──
            img_margin = Inches(1.0)
            img_area_left = img_margin
            img_area_top = Inches(1.15)
            img_area_width = self.sw - img_margin * 2
            img_area_height = self.sh - Inches(2.5)

            _add_rounded_rect(
                slide,
                left=img_area_left - Inches(0.1),
                top=img_area_top - Inches(0.1),
                width=img_area_width + Inches(0.2),
                height=img_area_height + Inches(0.2),
                fill_color=t["white"],
                line_color=t["border_gray"],
                line_width=1,
            )

            self._place_image_in_area(
                slide, slide_data, images_dir,
                img_area_left, img_area_top, img_area_width, img_area_height, t
            )

            # 图注放在图片正下方（居中）
            caption = slide_data.get("image_caption", "")
            if caption:
                _add_textbox(
                    slide,
                    left=Inches(1.0),
                    top=self.sh - Inches(1.3),
                    width=self.sw - Inches(2.0),
                    height=Inches(0.5),
                    text=caption,
                    font_name=t["font_caption"],
                    size_pt=t["caption_size_pt"],
                    color=t["secondary_text"],
                    italic=True,
                    alignment=PP_ALIGN.CENTER,
                )

        self._add_corner_decorations(slide)
        self._add_bottom_bar(slide, page_num)
        _add_notes(slide, slide_data.get("notes", ""))

    def _place_image_in_area(self, slide, slide_data, images_dir,
                             area_left, area_top, area_width, area_height, t):
        """在指定区域内加载并居中放置图片。"""
        image_path = slide_data.get("image_path", "")
        actual_image_file = None
        if image_path:
            if os.path.isabs(image_path):
                actual_image_file = image_path
            elif images_dir:
                actual_image_file = os.path.join(images_dir, image_path)

        if actual_image_file and os.path.isfile(actual_image_file):
            try:
                from PIL import Image as PILImage
                pil_img = PILImage.open(actual_image_file)
                img_w, img_h = pil_img.size
                aspect = img_w / img_h

                avail_w = area_width - Inches(0.4)
                avail_h = area_height - Inches(0.4)
                if avail_w / avail_h > aspect:
                    final_h = avail_h
                    final_w = int(final_h * aspect)
                else:
                    final_w = avail_w
                    final_h = int(final_w / aspect)

                img_left = area_left + (area_width - final_w) // 2
                img_top = area_top + (area_height - final_h) // 2

                slide.shapes.add_picture(
                    actual_image_file,
                    img_left, img_top,
                    final_w, final_h,
                )
            except ImportError:
                slide.shapes.add_picture(
                    actual_image_file,
                    area_left + Inches(0.2),
                    area_top + Inches(0.2),
                    area_width - Inches(0.4),
                    area_height - Inches(0.4),
                )
            except Exception as e:
                print(f"警告: 无法加载图片 '{actual_image_file}': {e}")
                _add_textbox(
                    slide,
                    left=area_left,
                    top=area_top + area_height // 2 - Inches(0.25),
                    width=area_width,
                    height=Inches(0.5),
                    text="[图片加载失败]",
                    font_name=t["font_body"],
                    size_pt=16,
                    color=t["secondary_text"],
                    alignment=PP_ALIGN.CENTER,
                )
        else:
            _add_rounded_rect(
                slide,
                left=area_left + Inches(0.5),
                top=area_top + Inches(0.5),
                width=area_width - Inches(1.0),
                height=area_height - Inches(1.0),
                fill_color=t["highlight"],
                line_color=t["border_gray"],
                line_width=1,
            )
            _add_textbox(
                slide,
                left=area_left + Inches(1.0),
                top=area_top + area_height // 2 - Inches(0.25),
                width=area_width - Inches(2.0),
                height=Inches(0.5),
                text="Figure Placeholder",
                font_name=t["font_body"],
                size_pt=16,
                color=t["secondary_text"],
                alignment=PP_ALIGN.CENTER,
            )

    def build_conclusion(self, prs, slide_data, page_num):
        """构建结论页。"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = self.theme

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = t["background"]

        # 标题带（替代原 _add_top_bar + 单独标题）
        self._add_title_header(slide, slide_data.get("title", "总结"), t,
                               kicker=slide_data.get("kicker"))

        # 核心信息 — 更大的结论框样式
        key_message = slide_data.get("key_message", slide_data.get("title", ""))
        box_w = Inches(10.0)
        box_h = Inches(1.8)
        box_left = (self.sw - box_w) // 2
        box_top = Inches(1.4)

        # 浅蓝色圆角矩形背景
        _add_rounded_rect(
            slide,
            left=box_left,
            top=box_top,
            width=box_w,
            height=box_h,
            fill_color=t["conclusion_bg"],
            line_color=None,
        )
        # 左侧蓝色竖线装饰（更粗）
        _add_rect(
            slide,
            left=box_left,
            top=box_top,
            width=Inches(0.1),
            height=box_h,
            fill_color=t["accent"],
        )

        # 核心信息文字（"结论：" 前缀 + 正文，更大字号，支持关键词高亮）
        kbox = slide.shapes.add_textbox(
            box_left + Inches(0.4),
            box_top + Inches(0.15),
            box_w - Inches(0.7),
            box_h - Inches(0.3),
        )
        ktf = kbox.text_frame
        ktf.word_wrap = True
        try:
            ktf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        kp = ktf.paragraphs[0]
        kp.alignment = PP_ALIGN.LEFT
        # "结论：" 前缀（加粗）
        prefix_run = kp.add_run()
        prefix_run.text = "结论："
        _set_font(prefix_run, t["font_title"], 22, bold=True,
                  color=t["header_bar"])
        # 正文（更大字号 + 关键词高亮）
        self._add_keyword_runs(kp, key_message, t["font_title"], 22,
                               t["header_bar"], t, bold=True)

        # 支撑要点
        bullets = slide_data.get("bullets", [])
        y = Inches(3.7)
        for i, bullet in enumerate(bullets):
            # 小装饰条
            _add_rect(
                slide,
                left=Inches(2.5),
                top=y + Inches(0.08),
                width=Inches(0.06),
                height=Inches(0.25),
                fill_color=t["accent"],
            )
            # 要点文本（支持关键词高亮）
            self._add_rich_textbox(
                slide,
                left=Inches(2.8),
                top=y,
                width=Inches(8.0),
                height=Inches(0.45),
                text=bullet,
                font_name=t["font_body"],
                size_pt=t["body_size_pt"],
                color=t["body_text"],
                theme=t,
            )
            y += Inches(0.6)

        self._add_corner_decorations(slide)
        self._add_bottom_bar(slide, page_num)
        _add_notes(slide, slide_data.get("notes", ""))

    def build_qa(self, prs, slide_data, page_num):
        """构建问答/致谢页。"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = self.theme

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = t["background"]

        # 标题带（替代原 _add_top_bar + 单独标题）
        self._add_title_header(slide, slide_data.get("title", "Thank You"), t,
                               kicker=slide_data.get("kicker"))

        # 中央大文字 "Q & A"（使用 key_message 或默认值）
        key_message = slide_data.get("key_message", "Q & A")
        _add_textbox(
            slide,
            left=Inches(2.0),
            top=Inches(1.8),
            width=self.sw - Inches(4.0),
            height=Inches(1.5),
            text=key_message,
            font_name=t["font_title"],
            size_pt=48,
            color=t["title_band"],
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # 装饰分隔线
        _add_rect(
            slide,
            left=Inches(4.5),
            top=Inches(3.5),
            width=Inches(4.33),
            height=Inches(0.04),
            fill_color=t["accent"],
        )

        # "Thank You" 文字
        _add_textbox(
            slide,
            left=Inches(2.0),
            top=Inches(3.8),
            width=self.sw - Inches(4.0),
            height=Inches(0.8),
            text="Thank You",
            font_name=t["font_title"],
            size_pt=32,
            color=t["accent"],
            alignment=PP_ALIGN.CENTER,
        )

        # 底部装饰圆
        if t["decorations"]:
            c1 = _add_oval(
                slide,
                left=Inches(2.5),
                top=Inches(5.0),
                width=Inches(0.8),
                height=Inches(0.8),
                fill_color=t["decor_corner"],
            )
            _set_shape_opacity(c1, 40)
            c2 = _add_oval(
                slide,
                left=self.sw - Inches(3.5),
                top=Inches(5.0),
                width=Inches(0.8),
                height=Inches(0.8),
                fill_color=t["decor_corner"],
            )
            _set_shape_opacity(c2, 40)
            c3 = _add_oval(
                slide,
                left=Inches(5.8),
                top=Inches(5.3),
                width=Inches(1.2),
                height=Inches(1.2),
                fill_color=t["accent"],
            )
            _set_shape_opacity(c3, 20)

        # 副信息
        bullets = slide_data.get("bullets", [])
        if bullets:
            y = Inches(4.8)
            for bullet in bullets:
                _add_textbox(
                    slide,
                    left=Inches(3.0),
                    top=y,
                    width=Inches(7.0),
                    height=Inches(0.4),
                    text=bullet,
                    font_name=t["font_body"],
                    size_pt=14,
                    color=t["secondary_text"],
                    alignment=PP_ALIGN.CENTER,
                )
                y += Inches(0.45)

        self._add_bottom_bar(slide, page_num)
        _add_notes(slide, slide_data.get("notes", ""))

    # ── ref 主题页面构建方法 ────────────────────────────

    def ref_build_cover(self, prs, slide_data, page_num):
        """ref 封面页：深蓝顶部粗条 + 居中白色直角卡片 + 深蓝标题 + 引文直角条。

        仿照用户提供的参考模板设计。
        封面页不显示导航栏。
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = self.theme

        self._ref_set_background(slide)

        # ── 顶部装饰条：深蓝色粗条 ──
        _add_rect(
            slide,
            left=Inches(0),
            top=Inches(0),
            width=self.sw,
            height=Inches(0.12),
            fill_color=t["accent3"],
        )

        # ── 居中白色直角卡片 ──
        container_w = Inches(10.5)
        container_h = Inches(4.5)
        container_left = (self.sw - container_w) // 2
        container_top = (self.sh - container_h) // 2 - Inches(0.2)
        _add_rect(
            slide,
            left=container_left,
            top=container_top,
            width=container_w,
            height=container_h,
            fill_color=t["white"],
            line_color=None,
        )

        # ── 主标题：优先中文标题 cn_title，否则回退到 title ──
        cn_title = slide_data.get("cn_title", "")
        en_title = slide_data.get("en_title", "")
        fallback_title = slide_data.get("title", "论文标题")
        main_title = cn_title if cn_title else fallback_title

        title_tb = _add_textbox(
            slide,
            left=container_left + Inches(0.6),
            top=container_top + Inches(0.7),
            width=container_w - Inches(1.2),
            height=Inches(1.2),
            text=main_title,
            font_name=t["font_title"],
            size_pt=26,
            color=t["accent3"],
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        # 缩小行距让多行标题更紧凑
        for p in title_tb.text_frame.paragraphs:
            p.space_after = Pt(2)
            p.space_before = Pt(2)

        # ── 英文副标题（居中，灰色） ──
        underline_top = container_top + Inches(2.1) if cn_title else container_top + Inches(2.15)
        # 优先 en_title，其次回退到 subtitle（当 cn_title 存在时）
        en_sub = ""
        if cn_title:
            en_sub = en_title if en_title else slide_data.get("subtitle", "")
        else:
            en_sub = slide_data.get("en_subtitle", "") or slide_data.get("subtitle", "")

        if en_sub:
            _add_textbox(
                slide,
                left=container_left + Inches(0.8),
                top=underline_top + Inches(0.2),
                width=container_w - Inches(1.6),
                height=Inches(0.5),
                text=en_sub,
                font_name=t["font_body"],
                size_pt=13,
                color=t["secondary_text"],
                alignment=PP_ALIGN.CENTER,
            )

        # ── 引文直角条（深蓝底白字，居中） ──
        # 优先使用 subtitle 作为引文，若有 bullets 则用 bullets
        cite_text = ""
        info_lines = slide_data.get("bullets", [])
        if info_lines:
            cite_text = "   |   ".join(info_lines)
        else:
            cite_text = slide_data.get("subtitle", "")

        if cite_text:
            pill_w = Inches(min(len(cite_text) * 0.09 + 0.8, 8.0))
            pill_h = Inches(0.38)
            pill_left = (self.sw - pill_w) // 2
            pill_top = container_top + container_h - Inches(0.75)
            _add_rect(
                slide,
                left=pill_left,
                top=pill_top,
                width=pill_w,
                height=pill_h,
                fill_color=t["accent3"],
                line_color=None,
            )
            _add_textbox(
                slide,
                left=pill_left + Inches(0.2),
                top=pill_top + Inches(0.02),
                width=pill_w - Inches(0.4),
                height=pill_h - Inches(0.04),
                text=cite_text,
                font_name=t["font_body"],
                size_pt=11,
                color=t["white"],
                bold=True,
                alignment=PP_ALIGN.CENTER,
            )

        # ── 底部信息（汇报人 + 日期） ──
        _add_textbox(
            slide,
            left=Inches(2.0),
            top=self.sh - Inches(0.85),
            width=self.sw - Inches(4.0),
            height=Inches(0.3),
            text=slide_data.get("presenter", "汇报人：研究生组会汇报"),
            font_name=t["font_body"],
            size_pt=10,
            color=t["tertiary_text"],
            alignment=PP_ALIGN.CENTER,
        )
        _add_textbox(
            slide,
            left=Inches(2.0),
            top=self.sh - Inches(0.55),
            width=self.sw - Inches(4.0),
            height=Inches(0.3),
            text=slide_data.get("date", "2026年7月29日"),
            font_name=t["font_body"],
            size_pt=10,
            color=t["tertiary_text"],
            alignment=PP_ALIGN.CENTER,
        )

        _add_notes(slide, slide_data.get("notes", ""))

    def ref_build_toc(self, prs, slide_data, page_num):
        """ref 目录页：导航栏 + 标题 + 编号卡片列表。

        每个目录项为白色圆角矩形卡片，左侧带 accent 竖条，内含编号和章节名。
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = self.theme

        self._ref_set_background(slide)
        self._ref_add_nav_bar(slide, page_num)
        self._ref_add_title(slide, slide_data.get("title", "目录"))

        sections = slide_data.get("sections", [])
        highlights = slide_data.get("highlights", [])

        # 卡片布局：两列布局，每列宽度约 5.8"
        col_gap = Inches(0.4)
        card_w = Inches(5.8)
        card_h = Inches(0.85)
        margin = Inches(t["margin_inch"])
        start_y = Inches(1.85)
        row_gap = Inches(0.25)

        for i, section in enumerate(sections):
            col = i % 2
            row = i // 2
            left = margin + col * (card_w + col_gap)
            top = start_y + row * (card_h + row_gap)

            # 编号卡片（白色圆角矩形 + accent 竖条）
            self._ref_add_card(slide, left, top, card_w, card_h,
                               accent_color=t["accent"])

            # 编号圆形（accent3 深蓝色）
            circle_size = Inches(0.42)
            circle_left = left + Inches(0.25)
            circle_top = top + (card_h - circle_size) // 2
            circle = _add_oval(
                slide,
                left=circle_left,
                top=circle_top,
                width=circle_size,
                height=circle_size,
                fill_color=t["accent3"],
            )
            ctf = circle.text_frame
            ctf.word_wrap = False
            try:
                ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
            cp = ctf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            crun = cp.add_run()
            crun.text = str(i + 1)
            _set_font(crun, t["font_body"], 14, bold=True, color=t["white"])

            # 章节名
            section_title = _extract_section_title(section)
            _add_textbox(
                slide,
                left=circle_left + circle_size + Inches(0.2),
                top=top + Inches(0.12),
                width=card_w - Inches(1.2),
                height=Inches(0.35),
                text=section_title,
                font_name=t["font_title"],
                size_pt=14,
                color=t["accent3"],
                bold=True,
            )
            # 章节描述（如有 highlights）
            if i < len(highlights):
                desc = highlights[i]
                if isinstance(desc, dict):
                    desc = desc.get("content", desc.get("title", ""))
                _add_textbox(
                    slide,
                    left=circle_left + circle_size + Inches(0.2),
                    top=top + Inches(0.48),
                    width=card_w - Inches(1.2),
                    height=Inches(0.3),
                    text=str(desc),
                    font_name=t["font_body"],
                    size_pt=10,
                    color=t["secondary_text"],
                )

        self._ref_add_bottom_bar(slide, page_num)
        _add_notes(slide, slide_data.get("notes", ""))

    def ref_build_section(self, prs, slide_data, page_num):
        """ref 章节分隔页：导航栏 + 居中章节编号和标题。

        章节标题居中，上下带装饰横线。
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = self.theme

        self._ref_set_background(slide)
        self._ref_add_nav_bar(slide, page_num)

        # 居中章节编号（深蓝色 accent3）
        section_no = slide_data.get("section_no", page_num)
        _add_textbox(
            slide,
            left=Inches(2.0),
            top=Inches(2.4),
            width=self.sw - Inches(4.0),
            height=Inches(0.6),
            text=f"PART {section_no}",
            font_name=t["font_title"],
            size_pt=18,
            color=t["accent3"],
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # 上方装饰横线
        line_w = Inches(3.0)
        _add_rect(
            slide,
            left=(self.sw - line_w) // 2,
            top=Inches(3.15),
            width=line_w,
            height=Inches(0.04),
            fill_color=t["accent"],
        )

        # 章节名（居中大号）
        section_title = slide_data.get("title", "")
        _add_textbox(
            slide,
            left=Inches(1.5),
            top=Inches(3.4),
            width=self.sw - Inches(3.0),
            height=Inches(1.0),
            text=section_title,
            font_name=t["font_title"],
            size_pt=36,
            color=t["accent3"],
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # 下方装饰横线
        _add_rect(
            slide,
            left=(self.sw - line_w) // 2,
            top=Inches(4.55),
            width=line_w,
            height=Inches(0.04),
            fill_color=t["accent"],
        )

        self._ref_add_bottom_bar(slide, page_num)
        _add_notes(slide, slide_data.get("notes", ""))

    def ref_build_content(self, prs, slide_data, page_num):
        """ref 内容页：导航栏 + 标题 + 白色卡片（含小标题、文本行、结论框）。

        文本行无项目符号，每行独立文本框；底部可选结论框。
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = self.theme

        self._ref_set_background(slide)
        self._ref_add_nav_bar(slide, page_num)
        self._ref_add_title(slide, slide_data.get("title", ""))

        margin = Inches(t["margin_inch"])
        # 主内容卡片（白色圆角矩形 + accent 竖条）
        card_left = margin
        card_top = Inches(1.85)
        card_w = self.sw - margin * 2
        card_h = self.sh - card_top - Inches(t["bottom_bar_height_inch"]) - Inches(0.25)
        self._ref_add_card(slide, card_left, card_top, card_w, card_h,
                           accent_color=t["accent"])

        # 卡片内边距
        inner_left = card_left + Inches(0.3)
        inner_top = card_top + Inches(0.25)
        inner_w = card_w - Inches(0.6)

        y = inner_top

        # 小标题（sub_title 或 analysis_title）
        sub_title = slide_data.get("sub_title") or slide_data.get("analysis_title")
        if sub_title:
            # 小标题（accent3 色加粗，无左侧竖条）
            _add_textbox(
                slide,
                left=inner_left,
                top=y,
                width=inner_w,
                height=Inches(0.34),
                text=sub_title,
                font_name=t["font_title"],
                size_pt=14,
                color=t["accent3"],
                bold=True,
            )
            y += Inches(0.5)

        # 文本行（无项目符号）
        bullets = slide_data.get("bullets", [])
        if bullets:
            max_lines = min(len(bullets), 10)
            line_y = self._ref_add_text_lines(
                slide,
                left=inner_left + Inches(0.15),
                top=y,
                width=inner_w - Inches(0.3),
                lines=bullets[:max_lines],
                font_size=t["body_size_pt"],
                color=t["body_text"],
                line_height=0.48,
            )
            y = line_y + Inches(0.15)

        # 结论框（如果有 conclusion 或 key_message）
        conclusion = slide_data.get("conclusion")
        if not conclusion:
            conclusion = slide_data.get("key_message")
        if conclusion:
            self._ref_add_conclusion_box(
                slide,
                left=inner_left,
                top=y,
                width=inner_w,
                text=conclusion,
            )

        self._ref_add_bottom_bar(slide, page_num)
        _add_notes(slide, slide_data.get("notes", ""))

    def ref_build_figure(self, prs, slide_data, page_num, images_dir=None):
        """ref 图表页：导航栏 + 标题 + 左侧文字卡片 + 右侧图片卡片。

        左侧卡片：白色圆角矩形 + accent 竖条，内含小标题、文本行、结论框。
        右侧卡片：白色圆角矩形，内含图片和图注。
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = self.theme

        self._ref_set_background(slide)
        self._ref_add_nav_bar(slide, page_num)
        self._ref_add_title(slide, slide_data.get("title", "图表"))

        margin = Inches(t["margin_inch"])
        # 内容区域顶部和底部
        content_top = Inches(1.85)
        content_bottom = self.sh - Inches(t["bottom_bar_height_inch"]) - Inches(0.25)
        content_h = content_bottom - content_top

        # 左侧文字卡片：约 5.0" 宽
        left_card_w = Inches(5.0)
        left_card_left = margin
        # 右侧图片卡片：约 6.875" 宽，从约 6.0" 处开始
        col_gap = Inches(0.25)
        right_card_w = self.sw - margin - (left_card_left + left_card_w + col_gap)
        right_card_left = left_card_left + left_card_w + col_gap

        # ── 左侧文字卡片 ──
        self._ref_add_card(slide, left_card_left, content_top,
                           left_card_w, content_h, accent_color=t["accent"])

        inner_left = left_card_left + Inches(0.3)
        inner_w = left_card_w - Inches(0.6)
        y = content_top + Inches(0.25)

        # 小标题
        sub_title = slide_data.get("sub_title") or slide_data.get("analysis_title")
        if sub_title:
            _add_textbox(
                slide,
                left=inner_left,
                top=y,
                width=inner_w,
                height=Inches(0.34),
                text=sub_title,
                font_name=t["font_title"],
                size_pt=13,
                color=t["accent3"],
                bold=True,
            )
            y += Inches(0.5)

        # 文本行
        bullets = slide_data.get("bullets", [])
        if bullets:
            max_lines = min(len(bullets), 8)
            line_y = self._ref_add_text_lines(
                slide,
                left=inner_left + Inches(0.15),
                top=y,
                width=inner_w - Inches(0.3),
                lines=bullets[:max_lines],
                font_size=t["body_size_pt"],
                color=t["body_text"],
                line_height=0.48,
            )
            y = line_y + Inches(0.15)

        # 结论框（贴在卡片底部）
        conclusion = slide_data.get("conclusion")
        if not conclusion:
            conclusion = slide_data.get("key_message")
        if conclusion:
            concl_top = content_bottom - Inches(0.75)
            self._ref_add_conclusion_box(
                slide,
                left=inner_left,
                top=concl_top,
                width=inner_w,
                text=conclusion,
            )

        # ── 右侧图片卡片 ──
        self._ref_add_card(slide, right_card_left, content_top,
                           right_card_w, content_h, accent_color=t["accent"])

        # 图片区域（卡片内边距）
        img_area_left = right_card_left + Inches(0.25)
        img_area_top = content_top + Inches(0.25)
        img_area_w = right_card_w - Inches(0.5)
        # 为图注预留底部空间
        img_area_h = content_h - Inches(0.7)

        self._place_image_in_area(
            slide, slide_data, images_dir,
            img_area_left, img_area_top, img_area_w, img_area_h, t
        )

        # 图注（图片正下方，居中）
        caption = slide_data.get("image_caption", "")
        if caption:
            _add_textbox(
                slide,
                left=img_area_left,
                top=img_area_top + img_area_h + Inches(0.05),
                width=img_area_w,
                height=Inches(0.3),
                text=caption,
                font_name=t["font_caption"],
                size_pt=t["caption_size_pt"],
                color=t["secondary_text"],
                italic=True,
                alignment=PP_ALIGN.CENTER,
            )

        self._ref_add_bottom_bar(slide, page_num)
        _add_notes(slide, slide_data.get("notes", ""))

    def ref_build_conclusion(self, prs, slide_data, page_num):
        """ref 结论页：导航栏 + 标题 + 核心信息卡片 + 支撑要点卡片。"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = self.theme

        self._ref_set_background(slide)
        self._ref_add_nav_bar(slide, page_num)
        self._ref_add_title(slide, slide_data.get("title", "总结"))

        margin = Inches(t["margin_inch"])
        content_top = Inches(1.85)
        content_bottom = self.sh - Inches(t["bottom_bar_height_inch"]) - Inches(0.25)

        # 核心信息卡片（顶部，accent 深蓝色竖条强调）
        key_message = slide_data.get("key_message", slide_data.get("title", ""))
        key_card_h = Inches(1.4)
        self._ref_add_card(slide, margin, content_top,
                           self.sw - margin * 2, key_card_h,
                           accent_color=t["accent"])

        # 核心信息文字（"结论：" 前缀 + 正文，支持关键词高亮）
        kbox = slide.shapes.add_textbox(
            margin + Inches(0.35),
            content_top + Inches(0.15),
            self.sw - margin * 2 - Inches(0.65),
            key_card_h - Inches(0.3),
        )
        ktf = kbox.text_frame
        ktf.word_wrap = True
        try:
            ktf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        kp = ktf.paragraphs[0]
        kp.alignment = PP_ALIGN.LEFT
        prefix_run = kp.add_run()
        prefix_run.text = "结论："
        _set_font(prefix_run, t["font_title"], 16, bold=True,
                  color=t["accent3"])
        self._add_keyword_runs(kp, key_message, t["font_title"], 16,
                               t["body_text"], t, bold=True)

        # 支撑要点（下方白色卡片）
        bullets = slide_data.get("bullets", [])
        if bullets:
            bullets_top = content_top + key_card_h + Inches(0.2)
            bullets_h = content_bottom - bullets_top
            self._ref_add_card(slide, margin, bullets_top,
                               self.sw - margin * 2, bullets_h,
                               accent_color=t["accent"])

            inner_left = margin + Inches(0.35)
            inner_w = self.sw - margin * 2 - Inches(0.7)
            y = bullets_top + Inches(0.2)

            max_bullets = min(len(bullets), 8)
            for i, bullet in enumerate(bullets[:max_bullets]):
                # 编号小方块（accent3 深蓝色）
                _add_rect(
                    slide,
                    left=inner_left,
                    top=y + Inches(0.04),
                    width=Inches(0.18),
                    height=Inches(0.18),
                    fill_color=t["accent3"],
                )
                # 要点文本（支持关键词高亮）
                self._add_rich_textbox(
                    slide,
                    left=inner_left + Inches(0.3),
                    top=y,
                    width=inner_w - Inches(0.3),
                    height=Inches(0.35),
                    text=bullet,
                    font_name=t["font_body"],
                    size_pt=t["body_size_pt"],
                    color=t["body_text"],
                    theme=t,
                )
                y += Inches(0.42)

        self._ref_add_bottom_bar(slide, page_num)
        _add_notes(slide, slide_data.get("notes", ""))

    def ref_build_qa(self, prs, slide_data, page_num):
        """ref 问答页：居中白色直角矩形容器 + "Q & A" 文字。

        问答页不显示导航栏。
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = self.theme

        self._ref_set_background(slide)

        # 居中白色直角矩形容器
        container_w = Inches(8.0)
        container_h = Inches(3.5)
        container_left = (self.sw - container_w) // 2
        container_top = (self.sh - container_h) // 2 - Inches(0.3)
        _add_rect(
            slide,
            left=container_left,
            top=container_top,
            width=container_w,
            height=container_h,
            fill_color=t["white"],
            line_color=None,
        )

        # "Q & A" 居中大字
        key_message = slide_data.get("key_message", "Q & A")
        _add_textbox(
            slide,
            left=container_left,
            top=container_top + Inches(1.0),
            width=container_w,
            height=Inches(1.5),
            text=key_message,
            font_name=t["font_title"],
            size_pt=48,
            color=t["accent3"],
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # 装饰横线（深蓝色 accent3）
        line_w = Inches(2.5)
        _add_rect(
            slide,
            left=(self.sw - line_w) // 2,
            top=container_top + Inches(2.4),
            width=line_w,
            height=Inches(0.04),
            fill_color=t["accent3"],
        )

        # 副信息（Thank You 等）
        _add_textbox(
            slide,
            left=container_left,
            top=container_top + Inches(2.6),
            width=container_w,
            height=Inches(0.5),
            text="Thank You",
            font_name=t["font_title"],
            size_pt=18,
            color=t["accent"],
            alignment=PP_ALIGN.CENTER,
        )

        # 底部副信息（如有 bullets）
        bullets = slide_data.get("bullets", [])
        if bullets:
            y = container_top + container_h + Inches(0.2)
            for bullet in bullets:
                _add_textbox(
                    slide,
                    left=Inches(3.0),
                    top=y,
                    width=self.sw - Inches(6.0),
                    height=Inches(0.3),
                    text=bullet,
                    font_name=t["font_body"],
                    size_pt=11,
                    color=t["secondary_text"],
                    alignment=PP_ALIGN.CENTER,
                )
                y += Inches(0.35)

        # 底部条（保持风格统一）
        self._ref_add_bottom_bar(slide, page_num)
        _add_notes(slide, slide_data.get("notes", ""))


# ═══════════════════════════════════════════════════════════
#  页面类型路由
# ═══════════════════════════════════════════════════════════

PAGE_TYPE_BUILDERS = {
    "cover":     "build_cover",
    "toc":       "build_toc",
    "section":   "build_section",
    "content":   "build_content",
    "figure":    "build_figure",
    "model":     "build_content",       # model 页复用 content 布局
    "conclusion":"build_conclusion",
    "qa":        "build_qa",
}


def build_slide(prs, builder, slide_data, page_idx, images_dir):
    """根据页面类型路由到对应的构建方法。

    对于 ref 主题，使用 ref_* 系列方法；其它主题沿用原 PAGE_TYPE_BUILDERS 路由。
    """
    page_type = slide_data.get("page_type", "content")

    # ── ref 主题专用路由 ──
    if builder.theme_name == "ref":
        ref_method_map = {
            "cover":      "ref_build_cover",
            "toc":        "ref_build_toc",
            "section":    "ref_build_section",
            "content":    "ref_build_content",
            "model":      "ref_build_content",   # model 页复用 content 布局
            "figure":     "ref_build_figure",
            "conclusion": "ref_build_conclusion",
            "qa":         "ref_build_qa",
        }
        method_name = ref_method_map.get(page_type, "ref_build_content")
        method = getattr(builder, method_name)
        if method_name == "ref_build_figure":
            method(prs, slide_data, page_idx, images_dir=images_dir)
        else:
            method(prs, slide_data, page_idx)
        return

    # ── 原始路由（academic / minimal / trae） ──
    method_name = PAGE_TYPE_BUILDERS.get(page_type, "build_content")
    method = getattr(builder, method_name)

    # figure 类型额外传入 images_dir
    if method_name == "build_figure":
        method(prs, slide_data, page_idx, images_dir=images_dir)
    else:
        method(prs, slide_data, page_idx)


# ═══════════════════════════════════════════════════════════
#  输入验证
# ═══════════════════════════════════════════════════════════

REQUIRED_FIELDS = {"page_num", "page_type", "title"}
OPTIONAL_FIELDS = {
    "subtitle", "sections", "bullets", "image_path", "image_caption",
    "notes", "highlights", "key_message",
    "sub_title", "analysis_title", "conclusion", "kicker", "section_no",
    "cn_title", "en_title", "label", "presenter", "date",
}


def _extract_section_title(section):
    """从 sections 数组项中提取标题文字，兼容字符串和字典两种格式。

    - 字符串: "研究背景" → "研究背景"
    - 字典:   {"title": "研究背景", "en": "Background", "desc": "..."} → "研究背景"
    """
    if isinstance(section, str):
        return section
    if isinstance(section, dict):
        return section.get("title", section.get("name", str(section)))
    return str(section)


def validate_slides(slides):
    """验证 slides.json 数据结构，返回 (errors, warnings)。

    errors:   阻断生成的错误（必需字段缺失、类型错误等）
    warnings: 跨环境一致性警告（R1-R7），不阻断生成但建议修复
    """
    errors = []
    warnings = []

    # ── R1: 固定17页结构 ──
    if len(slides) != 17:
        warnings.append(f"  [R1] 总页数应为17页，实际 {len(slides)} 页")

    for i, s in enumerate(slides):
        keys = set(s.keys())
        pt = s.get("page_type", "")

        # 检查必需字段（封面页可用 cn_title 替代 title）
        missing = REQUIRED_FIELDS - keys
        if pt == "cover" and "cn_title" in keys:
            missing.discard("title")
        if missing:
            errors.append(f"  幻灯片 #{i+1}: 缺少必需字段 {missing}")

        # 检查 page_type
        if "page_type" in s and s["page_type"] not in PAGE_TYPE_BUILDERS:
            errors.append(
                f"  幻灯片 #{i+1}: 未知 page_type '{s['page_type']}', "
                f"支持: {list(PAGE_TYPE_BUILDERS.keys())}"
            )

        # 检查 sections 格式（toc 页）
        sections = s.get("sections")
        if sections is not None:
            if not isinstance(sections, list):
                errors.append(f"  幻灯片 #{i+1}: sections 必须是数组")
            else:
                for j, sec in enumerate(sections):
                    if not isinstance(sec, (str, dict)):
                        errors.append(
                            f"  幻灯片 #{i+1}: sections[{j}] 必须是字符串或字典, "
                            f"实际类型: {type(sec).__name__}"
                        )
                # R5: toc 四段式
                if pt == "toc" and len(sections) != 4:
                    warnings.append(f"  [R5] 幻灯片 #{i+1}: toc 页 sections 应为4段，实际 {len(sections)} 段")

        # 检查 highlights 格式
        highlights = s.get("highlights")
        if highlights is not None:
            if not isinstance(highlights, list):
                errors.append(f"  幻灯片 #{i+1}: highlights 必须是数组")

        # R2: 封面 cn_title + en_title
        if pt == "cover":
            if "cn_title" not in keys:
                warnings.append(f"  [R2] 幻灯片 #{i+1}: 封面页缺少 cn_title（禁止仅用 title）")
            if "en_title" not in keys:
                warnings.append(f"  [R2] 幻灯片 #{i+1}: 封面页缺少 en_title")

        # R3/R4/R7: figure 页检查
        if pt == "figure":
            ip = s.get("image_path")
            if not ip:
                warnings.append(f"  [R3] 幻灯片 #{i+1}: figure 页缺少 image_path")
            elif isinstance(ip, list):
                warnings.append(f"  [R4] 幻灯片 #{i+1}: image_path 是数组，违反一图一页")
            if not s.get("image_caption"):
                warnings.append(f"  [R7] 幻灯片 #{i+1}: figure 页缺少 image_caption")
            if not s.get("sub_title") and not s.get("analysis_title"):
                warnings.append(f"  [R7] 幻灯片 #{i+1}: figure 页缺少 sub_title")

        # R7: content 页必填字段
        if pt == "content":
            if not s.get("sub_title") and not s.get("analysis_title"):
                warnings.append(f"  [R7] 幻灯片 #{i+1}: content 页缺少 sub_title")
            if not s.get("conclusion"):
                warnings.append(f"  [R7] 幻灯片 #{i+1}: content 页缺少 conclusion")

        # R6: bullets 数量和长度
        _bullet_limits = {"content": (3, 4, 40), "figure": (2, 3, 40), "conclusion": (3, 5, 35)}
        if pt in _bullet_limits:
            min_n, max_n, max_len = _bullet_limits[pt]
            bullets = s.get("bullets", [])
            if bullets:
                if len(bullets) < min_n or len(bullets) > max_n:
                    warnings.append(
                        f"  [R6] 幻灯片 #{i+1}: {pt} 页 bullets 应为 {min_n}-{max_n} 条，"
                        f"实际 {len(bullets)} 条"
                    )
                for j, b in enumerate(bullets):
                    if isinstance(b, str) and len(b) > max_len:
                        warnings.append(f"  [R6] 幻灯片 #{i+1}: bullets[{j}] 长度 {len(b)} 超过 {max_len}")

    # R1: page_type 序列检查
    actual_types = [s.get("page_type", "unknown") for s in slides]
    fig_count = actual_types.count("figure")
    expected_seq = ["cover", "toc", "section", "content", "content", "section"]
    expected_seq.extend(["figure"] * fig_count)
    expected_seq.extend(["section", "content", "conclusion", "qa"])
    for j in range(max(len(actual_types), len(expected_seq))):
        a = actual_types[j] if j < len(actual_types) else "(缺失)"
        e = expected_seq[j] if j < len(expected_seq) else "(多余)"
        if a != e:
            warnings.append(f"  [R1] 第{j+1}页 page_type 应为 '{e}'，实际 '{a}'")

    return errors, warnings


# ═══════════════════════════════════════════════════════════
#  主函数
# ═══════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description="gen_pptx —— paper-report-ppt 直接生成路径, "
                    "使用 python-pptx 生成高质量可编辑 PPTX。",
    )
    parser.add_argument(
        "--input", "-i",
        required=True,
        help="slides.json 路径 (AI 生成的幻灯片定义文件)",
    )
    parser.add_argument(
        "--images-dir",
        default="./images",
        help="图片目录路径 (默认: ./images)",
    )
    parser.add_argument(
        "--output", "-o",
        default="./output.pptx",
        help="输出 PPTX 文件路径 (默认: ./output.pptx)",
    )
    parser.add_argument(
        "--theme",
        default="ref",
        choices=["ref"],
        help="主题风格：ref=深蓝导航栏+白色直角卡片（默认，对齐参考模板，唯一可选）",
    )
    args = parser.parse_args()

    # ── 检查输入文件 ──
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"错误: 输入文件不存在: {args.input}")
        print("请确保 slides.json 文件存在，格式参见文档。")
        sys.exit(1)

    # ── 读取 JSON ──
    try:
        with open(input_path, "r", encoding="utf-8") as f:
            slides_data = json.load(f)
    except json.JSONDecodeError as e:
        print(f"错误: slides.json 格式不合法: {e}")
        sys.exit(1)

    if not isinstance(slides_data, list) or len(slides_data) == 0:
        print("错误: slides.json 应为非空数组 (list of slide objects)。")
        sys.exit(1)

    # ── 验证数据 ──
    errors, warnings = validate_slides(slides_data)
    if warnings:
        print("⚠️ 跨环境一致性警告（不阻止生成，但建议修复以确保跨环境一致）：")
        for warn in warnings:
            print(warn)
        print()
    if errors:
        print("数据验证错误:")
        for err in errors:
            print(err)
        sys.exit(1)

    # ── 预处理：为 section 页分配章节序号 ──
    section_counter = 0
    for slide_data in slides_data:
        if slide_data.get("page_type") == "section":
            section_counter += 1
            # 如果 slides.json 未提供 section_no，自动分配
            if "section_no" not in slide_data:
                slide_data["section_no"] = section_counter

    # ── 自动补全：确保每页都有足够的字段，效果不依赖AI生成质量 ──
    for slide_data in slides_data:
        page_type = slide_data.get("page_type", "content")
        # content/figure 页：如果没有 conclusion 但有 highlights，提取最后一条作为结论
        if page_type in ("content", "figure"):
            if "conclusion" not in slide_data:
                highlights = slide_data.get("highlights", [])
                if highlights:
                    if isinstance(highlights[-1], dict):
                        slide_data["conclusion"] = highlights[-1].get("content", "")
                    elif isinstance(highlights[-1], str):
                        slide_data["conclusion"] = highlights[-1]
                elif "key_message" in slide_data:
                    slide_data["conclusion"] = slide_data["key_message"]
        # conclusion 页：确保有 key_message
        if page_type == "conclusion" and "key_message" not in slide_data:
            bullets = slide_data.get("bullets", [])
            if bullets:
                slide_data["key_message"] = bullets[0] if isinstance(bullets[0], str) else str(bullets[0])

    # ── 构建演示文稿 ──
    builder = SlideBuilder(theme_name=args.theme)
    prs = builder.create_presentation()

    print(f"主题: {args.theme}")
    print(f"幻灯片数量: {len(slides_data)}")
    print(f"图片目录: {args.images_dir}")
    print("-" * 50)

    for i, slide_data in enumerate(slides_data):
        page_num = slide_data.get("page_num", i + 1)
        page_type = slide_data.get("page_type", "content")
        title = slide_data.get("title", "(无标题)")

        build_slide(prs, builder, slide_data, page_num, args.images_dir)

        print(f"  [{page_num}/{len(slides_data)}] {page_type}: {title}")

    # ── 保存 ──
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print("-" * 50)
    print(f"已生成: {output_path.resolve()}")
    print("完成!")


if __name__ == "__main__":
    main()
